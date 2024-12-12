"""Export manager for proposal generation."""
import os
from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import pptx
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.units import inch
from jinja2 import Environment, FileSystemLoader
import json
import shutil

class ExportManager:
    """Manages export operations for proposals."""

    def __init__(self, output_dir: str = "exports"):
        """Initialize export manager."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.template_dir = Path(__file__).parent / "templates"

    def export_pdf(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as PDF using ReportLab."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = self.output_dir / f"proposal_{timestamp}.pdf"
        
        # Create PDF document
        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=18,
            spaceAfter=20
        )
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=12
        )
        
        # Build content
        story = []
        
        # Title
        story.append(Paragraph(f"Project Proposal for {proposal_data['client_name']}", title_style))
        story.append(Paragraph(f"Generated on {proposal_data['generation_date']}", styles['Italic']))
        story.append(Spacer(1, 30))
        
        # Project Overview
        story.append(Paragraph("Project Overview", heading_style))
        story.append(Paragraph(proposal_data['project_description'], body_style))
        story.append(Spacer(1, 20))
        
        # Industry & Target Market
        data = [
            ["Industry", proposal_data['industry']],
            ["Target Market", proposal_data['target_market']]
        ]
        table = Table(data, colWidths=[2*inch, 4*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(table)
        story.append(Spacer(1, 30))
        
        # SEO Analysis
        if proposal_data.get('seo_insights'):
            story.append(Paragraph("SEO Analysis", heading_style))
            for key, value in proposal_data['seo_insights'].items():
                if isinstance(value, dict):
                    for subkey, subvalue in value.items():
                        story.append(Paragraph(f"<b>{key} - {subkey}:</b> {subvalue}", body_style))
                else:
                    story.append(Paragraph(f"<b>{key}:</b> {value}", body_style))
            story.append(Spacer(1, 20))
        
        # Market Analysis
        if proposal_data.get('market_insights'):
            story.append(Paragraph("Market Analysis", heading_style))
            for key, value in proposal_data['market_insights'].items():
                if isinstance(value, dict):
                    for subkey, subvalue in value.items():
                        story.append(Paragraph(f"<b>{key} - {subkey}:</b> {subvalue}", body_style))
                else:
                    story.append(Paragraph(f"<b>{key}:</b> {value}", body_style))
            story.append(Spacer(1, 20))
        
        # Content Strategy
        if proposal_data.get('content_plan'):
            story.append(Paragraph("Content Strategy", heading_style))
            for key, value in proposal_data['content_plan'].items():
                if isinstance(value, (list, tuple)):
                    story.append(Paragraph(f"<b>{key}:</b>", body_style))
                    for item in value:
                        story.append(Paragraph(f"• {item}", body_style))
                else:
                    story.append(Paragraph(f"<b>{key}:</b> {value}", body_style))
            story.append(Spacer(1, 20))
        
        # Mockups
        if proposal_data.get('mockups'):
            story.append(Paragraph("Design Mockups", heading_style))
            for device, image_path in proposal_data['mockups'].items():
                if os.path.exists(image_path):
                    story.append(Paragraph(f"{device.title()} View", body_style))
                    story.append(Image(image_path, width=6*inch, height=4*inch))
                    story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        return str(output_path)

    def export_pptx(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as PowerPoint presentation."""
        prs = pptx.Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = f"Proposal for {proposal_data['client_name']}"
        subtitle.text = f"Generated on {proposal_data['generation_date']}"
        
        # Overview slide
        overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = overview_slide.shapes.title
        content = overview_slide.placeholders[1]
        title.text = "Project Overview"
        content.text = proposal_data['project_description']
        
        # Add slides for each section
        self._add_section_slides(prs, "SEO Analysis", proposal_data.get('seo_insights', {}))
        self._add_section_slides(prs, "Market Analysis", proposal_data.get('market_insights', {}))
        self._add_section_slides(prs, "Content Strategy", proposal_data.get('content_plan', {}))
        
        # Add mockup slides
        if proposal_data.get('mockups'):
            self._add_mockup_slides(prs, proposal_data['mockups'])
        
        # Save presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = self.output_dir / f"proposal_{timestamp}.pptx"
        prs.save(str(output_path))
        
        return str(output_path)

    def export_json(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as JSON."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = self.output_dir / f"proposal_{timestamp}.json"
        
        with open(output_path, 'w') as f:
            json.dump(proposal_data, f, indent=2)
        
        return str(output_path)

    def export_archive(self, proposal_data: Dict[str, Any]) -> str:
        """Export all formats in a zip archive."""
        # Create temporary directory for exports
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        temp_dir = self.output_dir / f"temp_{timestamp}"
        temp_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            # Generate all formats
            pdf_path = Path(self.export_pdf(proposal_data))
            pptx_path = Path(self.export_pptx(proposal_data))
            json_path = Path(self.export_json(proposal_data))
            
            # Copy to temp directory
            shutil.copy2(pdf_path, temp_dir)
            shutil.copy2(pptx_path, temp_dir)
            shutil.copy2(json_path, temp_dir)
            
            # Create zip archive
            archive_path = self.output_dir / f"proposal_{timestamp}_complete.zip"
            shutil.make_archive(str(archive_path.with_suffix('')), 'zip', temp_dir)
            
            return str(archive_path)
        finally:
            # Cleanup
            shutil.rmtree(temp_dir)

    def _add_section_slides(self, prs: pptx.Presentation, section_title: str, section_data: Dict[str, Any]):
        """Add slides for a proposal section."""
        if not section_data:
            return
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = section_title
        content_text = []
        
        for key, value in section_data.items():
            if isinstance(value, (str, int, float)):
                content_text.append(f"• {key}: {value}")
            elif isinstance(value, (list, tuple)):
                content_text.append(f"• {key}:")
                for item in value:
                    content_text.append(f"  - {item}")
        
        content.text = "\n".join(content_text)

    def _add_mockup_slides(self, prs: pptx.Presentation, mockups: Dict[str, str]):
        """Add slides for mockups."""
        for device, image_path in mockups.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"{device.title()} Mockup"
            
            # Add mockup image
            if os.path.exists(image_path):
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                slide.shapes.add_picture(image_path, left, top, width=width) 