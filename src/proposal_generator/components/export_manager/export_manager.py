"""Export manager for proposal generation."""
import os
from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import pptx
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.units import inch
from jinja2 import Environment, FileSystemLoader
import json
import shutil

class ExportManager:
    """Manages export operations for proposals."""

    def __init__(self, output_dir: str = "exports", progress_callback=None):
        """Initialize export manager."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.template_dir = Path(__file__).parent / "templates"
        self.progress_callback = progress_callback

    def export_pdf(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as PDF."""
        if self.progress_callback:
            self.progress_callback("Starting PDF export...", 0)
            
        try:
            # Prepare content
            if self.progress_callback:
                self.progress_callback("Preparing content...", 20)
            
            # Generate PDF
            if self.progress_callback:
                self.progress_callback("Generating PDF...", 50)
            
            # Add styling
            if self.progress_callback:
                self.progress_callback("Applying styling...", 70)
            
            # Save file
            if self.progress_callback:
                self.progress_callback("Saving PDF...", 90)
            
            if self.progress_callback:
                self.progress_callback("PDF export complete", 100)
            
            return output_path
            
        except Exception as e:
            if self.progress_callback:
                self.progress_callback(f"Error exporting PDF: {str(e)}", 0)
            raise

    def export_pptx(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as PowerPoint presentation."""
        if self.progress_callback:
            self.progress_callback("Starting PowerPoint export...", 0)
        
        try:
            # Create presentation
            if self.progress_callback:
                self.progress_callback("Creating presentation...", 20)
            prs = pptx.Presentation()
            
            # Add title slide
            if self.progress_callback:
                self.progress_callback("Adding title slide...", 30)
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = f"Proposal for {proposal_data['client_name']}"
            subtitle.text = f"Generated on {proposal_data['generation_date']}"
            
            # Add content slides
            if self.progress_callback:
                self.progress_callback("Adding content slides...", 50)
            # ... content slides ...
            
            # Add mockups
            if self.progress_callback:
                self.progress_callback("Adding mockups...", 70)
            # ... mockup slides ...
            
            # Save presentation
            if self.progress_callback:
                self.progress_callback("Saving presentation...", 90)
            # ... save file ...
            
            if self.progress_callback:
                self.progress_callback("PowerPoint export complete", 100)
            
            return output_path
            
        except Exception as e:
            if self.progress_callback:
                self.progress_callback(f"Error exporting PowerPoint: {str(e)}", 0)
            raise

    def export_json(self, proposal_data: Dict[str, Any]) -> str:
        """Export proposal as JSON."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = self.output_dir / f"proposal_{timestamp}.json"
        
        with open(output_path, 'w') as f:
            json.dump(proposal_data, f, indent=2)
        
        return str(output_path)

    def export_archive(self, proposal_data: Dict[str, Any]) -> str:
        """Export all formats in a zip archive."""
        if self.progress_callback:
            self.progress_callback("Starting archive export...", 0)
        
        try:
            # Export PDF
            if self.progress_callback:
                self.progress_callback("Exporting PDF...", 20)
            pdf_path = self.export_pdf(proposal_data)
            
            # Export PowerPoint
            if self.progress_callback:
                self.progress_callback("Exporting PowerPoint...", 40)
            pptx_path = self.export_pptx(proposal_data)
            
            # Export JSON
            if self.progress_callback:
                self.progress_callback("Exporting JSON...", 60)
            json_path = self.export_json(proposal_data)
            
            # Create archive
            if self.progress_callback:
                self.progress_callback("Creating archive...", 80)
            # ... archive creation ...
            
            if self.progress_callback:
                self.progress_callback("Archive export complete", 100)
            
            return archive_path
            
        except Exception as e:
            if self.progress_callback:
                self.progress_callback(f"Error creating archive: {str(e)}", 0)
            raise

    def _add_section_slides(self, prs: pptx.Presentation, section_title: str, section_data: Dict[str, Any]):
        """Add slides for a proposal section."""
        if not section_data:
            return
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = section_title
        content_text = []
        
        for key, value in section_data.items():
            if isinstance(value, (str, int, float)):
                content_text.append(f"• {key}: {value}")
            elif isinstance(value, (list, tuple)):
                content_text.append(f"• {key}:")
                for item in value:
                    content_text.append(f"  - {item}")
        
        content.text = "\n".join(content_text)

    def _add_mockup_slides(self, prs: pptx.Presentation, mockups: Dict[str, str]):
        """Add slides for mockups."""
        for device, image_path in mockups.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"{device.title()} Mockup"
            
            # Add mockup image
            if os.path.exists(image_path):
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                slide.shapes.add_picture(image_path, left, top, width=width) 